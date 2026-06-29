import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")
    doc = Document(str(path))
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        style = para.style.name.lower() if para.style else ""
        if "heading" in style:
            level = 1
            for s in ("heading 1", "heading 2", "heading 3", "heading 4", "heading 5", "heading 6"):
                if s in style:
                    level = int(s.split()[-1])
                    break
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)

    for table in doc.tables:
        lines.append("")
        for row in table.rows:
            cells = []
            seen = set()
            for cell in row.cells:
                cell_id = id(cell._tc)
                if cell_id in seen:
                    continue
                seen.add(cell_id)
                cells.append(cell.text.strip())
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")

    return "\n".join(lines)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(str(path))
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


def extract_html(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="replace")
    try:
        import html2text
        h = html2text.HTML2Text()
        h.body_width = 0
        h.ignore_links = False
        h.ignore_images = False
        h.ignore_tables = False
        return h.handle(content).strip()
    except ImportError:
        pass
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(content, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return text
    except ImportError:
        pass
    import re
    clean = re.sub(r"<[^>]+>", " ", content)
    clean = re.sub(r"\s+", " ", clean).strip()
    return clean


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## {sheet_name}")
        lines.append("")

        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    wb.close()
    return "\n".join(lines)


def extract_epub(path: Path) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("ebooklib and beautifulsoup4 required. Run: pip install ebooklib beautifulsoup4")

    book = epub.read_epub(str(path))
    lines = []

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            content = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(content, "html.parser")
            title_tag = soup.find("title")
            if title_tag and title_tag.get_text(strip=True):
                lines.append(f"## {title_tag.get_text(strip=True)}")
                lines.append("")

            for tag in soup(["script", "style", "nav"]):
                tag.decompose()

            for heading in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                level = int(heading.name[1])
                lines.append(f"{'#' * level} {heading.get_text(strip=True)}")

            for p in soup.find_all("p"):
                text = p.get_text(strip=True)
                if text:
                    lines.append(text)

            for img in soup.find_all("img"):
                alt = img.get("alt", "") or img.get("title", "") or "Imagem"
                src = img.get("src", "")
                lines.append(f"[Imagem: {alt}]")

    return "\n".join(lines) if lines else ""


def extract_csv(path: Path) -> str:
    import csv
    lines = []
    with open(str(path), "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def extract_markdown(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_latex(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("beautifulsoup4 not installed. Run: pip install beautifulsoup4")

    try:
        import subprocess
        result = subprocess.run(
            ["pandoc", str(path), "-f", "latex", "-t", "html"],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            soup = BeautifulSoup(result.stdout, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            return text
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception):
        pass

    import re
    content = re.sub(r'\\(?:begin|end)\{[^}]*\}', '', content)
    content = re.sub(r'\\[a-zA-Z]+(\{[^}]*\})*', '', content)
    content = re.sub(r'[{}]', '', content)
    content = re.sub(r'\n{3,}', '\n\n', content)
    return content.strip()


EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".epub": extract_epub,
    ".csv": extract_csv,
    ".md": extract_markdown,
    ".tex": extract_latex,
    ".txt": extract_markdown,
}


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported format: {ext}")
    logger.info("Extracting native text from %s file: %s", ext, path.name)
    return extractor(path)
